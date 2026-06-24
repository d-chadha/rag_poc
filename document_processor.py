"""
document_processor.py
---------------------
Extracts plain text from every supported document format and splits it into
chunks suitable for embedding.

Supported formats
-----------------
PDF   · DOCX  · PPTX  · XLSX  · TXT  · MD   · HTML · RTF  · CSV
"""

from __future__ import annotations

import csv
import io
import textwrap
from dataclasses import dataclass, field
from pathlib import Path
from typing import List

import chardet


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass
class DocumentChunk:
    """A single chunk of text ready to be embedded."""
    text: str
    metadata: dict = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets: List[str] = []
    for sheet in wb.worksheets:
        rows: List[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append("\t".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _extract_html(data: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(data, "html.parser")
    for tag in soup(["script", "style", "meta", "head"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extract_markdown(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    # Strip fenced code blocks (keep content), remove headers' '#' markers, etc.
    # For RAG we just want raw text — markdown syntax is low-signal noise.
    import re
    text = re.sub(r"```.*?```", lambda m: m.group(0).replace("```", ""), text, flags=re.DOTALL)
    text = re.sub(r"[#*`_~>\-]+", " ", text)
    return text


def _extract_rtf(data: bytes) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(data.decode("latin-1", errors="replace"))


def _extract_csv(data: bytes) -> str:
    encoding = chardet.detect(data).get("encoding") or "utf-8"
    text = data.decode(encoding, errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(row) for row in reader]
    return "\n".join(rows)


def _extract_txt(data: bytes) -> str:
    encoding = chardet.detect(data).get("encoding") or "utf-8"
    return data.decode(encoding, errors="replace")


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xlsx,
    ".html": _extract_html,
    ".htm":  _extract_html,
    ".md":   _extract_markdown,
    ".rtf":  _extract_rtf,
    ".csv":  _extract_csv,
    ".txt":  _extract_txt,
}

SUPPORTED_EXTENSIONS = sorted(_EXTRACTORS.keys())


def extract_text(filename: str, data: bytes) -> str:
    """Return plain text from *data* using the appropriate extractor."""
    ext = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )
    return extractor(data)


# ---------------------------------------------------------------------------
# Chunker
# ---------------------------------------------------------------------------

def chunk_text(
    text: str,
    filename: str,
    chunk_size: int = 800,
    chunk_overlap: int = 150,
) -> List[DocumentChunk]:
    """
    Split *text* into overlapping chunks.

    Parameters
    ----------
    chunk_size    : approximate number of characters per chunk
    chunk_overlap : characters shared between adjacent chunks
    """
    text = text.strip()
    if not text:
        return []

    chunks: List[DocumentChunk] = []
    start = 0
    total = len(text)
    chunk_index = 0

    while start < total:
        end = min(start + chunk_size, total)

        # Try to break at the last newline within the window so we don't
        # cut mid-sentence when possible.
        if end < total:
            last_nl = text.rfind("\n", start, end)
            if last_nl > start + chunk_overlap:
                end = last_nl + 1

        chunk_text_str = text[start:end].strip()
        if chunk_text_str:
            chunks.append(
                DocumentChunk(
                    text=chunk_text_str,
                    metadata={
                        "source": filename,
                        "chunk_index": chunk_index,
                        "start_char": start,
                        "end_char": end,
                    },
                )
            )
            chunk_index += 1

        start = end - chunk_overlap
        if start >= end:          # safety: never go backwards
            start = end

    return chunks


# ---------------------------------------------------------------------------
# High-level helper
# ---------------------------------------------------------------------------

def process_document(
    filename: str,
    data: bytes,
    chunk_size: int = 800,
    chunk_overlap: int = 150,
) -> List[DocumentChunk]:
    """Extract text from *data* and return a list of overlapping chunks."""
    text = extract_text(filename, data)
    return chunk_text(text, filename, chunk_size, chunk_overlap)
